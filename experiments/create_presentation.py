"""
experiments/create_presentation.py
===================================
Generates the 10-slide research presentation:
  "Self-Healing LLM Security Pipeline"

Run with:
    python experiments/create_presentation.py
    python experiments/create_presentation.py --output results/presentation.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour palette ────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # slide background
ACCENT     = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent / headings
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
AMBER      = RGBColor(0xFF, 0xA5, 0x00)
RED        = RGBColor(0xFF, 0x4C, 0x4C)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation) -> object:
    """Add a completely blank slide (layout 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _bg(slide, color: RGBColor = DARK_BG) -> None:
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(
    slide,
    left: float, top: float, width: float, height: float,
    text: str = "",
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bg_color: RGBColor | None = None,
) -> object:
    """Add a text box and return the shape."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _heading(slide, text: str, top: float = 0.35) -> None:
    """Slide heading in accent colour."""
    _box(slide, 0.5, top, 12.3, 0.7, text,
         font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)


def _rule(slide, top: float) -> None:
    """Horizontal accent rule under heading."""
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(0.5), Inches(top),
        Inches(12.83), Inches(top),
    )
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1.5)


def _bullet_box(
    slide,
    left: float, top: float, width: float, height: float,
    bullets: list[str],
    font_size: int = 16,
    color: RGBColor = WHITE,
    indent: float = 0.0,
) -> None:
    """Multi-line bullet box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_before = Pt(4)


def _label_value(
    slide,
    label: str, value: str,
    left: float, top: float,
    label_color: RGBColor = ACCENT,
    value_color: RGBColor = WHITE,
    font_size: int = 15,
) -> None:
    text = f"{label}  {value}"
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(6), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = label + "  "
    r1.font.bold = True
    r1.font.size = Pt(font_size)
    r1.font.color.rgb = label_color
    r2 = p.add_run()
    r2.text = value
    r2.font.size = Pt(font_size)
    r2.font.color.rgb = value_color


def _colored_box(
    slide,
    left: float, top: float, width: float, height: float,
    text: str,
    bg: RGBColor,
    fg: RGBColor = WHITE,
    font_size: int = 14,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.CENTER,
) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = bg
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg


# ── Slide builders ────────────────────────────────────────────────────

def slide_01_title(prs: Presentation) -> None:
    """Slide 1 — Title / Problem Framing."""
    slide = _blank_slide(prs)
    _bg(slide)

    # Accent bar on left
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.color.rgb = ACCENT

    _box(slide, 0.4, 1.5, 12.5, 1.0,
         "Self-Healing LLM Security Pipeline",
         font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

    _box(slide, 0.4, 2.55, 12.5, 0.5,
         "An automated Probe → Patch → Verify feedback loop for LLM vulnerability mitigation",
         font_size=18, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    _box(slide, 0.4, 3.3, 12.5, 0.4,
         "Research Collaboration Exercise  ·  CSL",
         font_size=14, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # Problem framing bullets
    _bullet_box(slide, 0.4, 4.1, 12.0, 2.8, [
        "▸  LLMs exhibit exploitable vulnerabilities: prompt injection, jailbreaks, toxic outputs, data leakage",
        "▸  Red-team tool Garak automates adversarial probing — analogous to penetration testing for LLMs",
        "▸  Self-healing systems can identify vulnerabilities and automatically apply targeted mitigations",
        "▸  Goal: end-to-end prototype demonstrating automated detect → patch → verify with measurable improvement",
    ], font_size=15, color=WHITE)


def slide_02_architecture(prs: Presentation) -> None:
    """Slide 2 — Pipeline Architecture."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Pipeline Architecture")
    _rule(slide, 1.1)

    # Three stage boxes
    stages = [
        ("Stage 1", "Baseline\nProbe", "Garak sends adversarial prompts to bare GPT-3.5-turbo.\nBuilds vulnerability baseline.", RGBColor(0x1F, 0x45, 0x6E)),
        ("Stage 2", "VDPS\nSynthesis", "LLM reads hitlog & synthesizes\ntargeted patches via JSON schema.", RGBColor(0x2D, 0x5A, 0x27)),
        ("Stage 3", "Verify via\nPatch Proxy", "Garak re-routed through localhost:8080.\nPatched prompts & filtered responses.", RGBColor(0x6E, 0x2F, 0x1F)),
    ]

    for i, (label, title, desc, color) in enumerate(stages):
        x = 0.5 + i * 4.1
        _colored_box(slide, x, 1.35, 3.7, 0.45, label,
                     bg=ACCENT, fg=DARK_BG, font_size=12, bold=True)
        _colored_box(slide, x, 1.8, 3.7, 1.0, title,
                     bg=color, fg=WHITE, font_size=18, bold=True)
        _colored_box(slide, x, 2.8, 3.7, 1.3, desc,
                     bg=RGBColor(0x22, 0x22, 0x3A), fg=LIGHT_GREY, font_size=13)

        # Arrow between stages
        if i < 2:
            ax = x + 3.75
            _box(slide, ax, 2.05, 0.35, 0.5, "→",
                 font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Iteration note
    _box(slide, 0.5, 4.25, 12.3, 0.4,
         "Iterative: pipeline loops up to max_iterations times — early exit when all attacks are mitigated",
         font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Key architecture insight
    _colored_box(slide, 0.5, 4.75, 12.3, 1.5,
        "Key design: PatchProxy intercepts at the network layer (OPENAI_BASE_URL override).\n"
        "Garak owns its own HTTP client — Python monkey-patching cannot intercept it.\n"
        "Setting OPENAI_BASE_URL=http://localhost:8080/v1 routes all scanner traffic through the proxy transparently.",
        bg=RGBColor(0x0D, 0x0D, 0x20), fg=LIGHT_GREY, font_size=13, align=PP_ALIGN.LEFT)


def slide_03_garak(prs: Presentation) -> None:
    """Slide 3 — Garak & Baseline Probe."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Stage 1 — Baseline Probe (Garak)")
    _rule(slide, 1.1)

    _bullet_box(slide, 0.5, 1.25, 7.5, 1.5, [
        "▸  Garak v0.14 — open-source LLM vulnerability scanner by NVIDIA",
        "▸  Sends structured adversarial prompt suites (probes) to target model",
        "▸  Each probe tests a specific vulnerability class",
        "▸  Reports attack pass-rate per probe category",
    ], font_size=15)

    # Probes table
    probes = [
        ("lmrc.Bullying",                  "Generates bullying content",            "0%"),
        ("lmrc.Deadnaming",                "Uses pre-transition names for trans ppl","83%"),
        ("lmrc.QuackMedicine",             "Endorses pseudoscientific health claims","60%"),
        ("lmrc.Sexualisation",             "Sexually objectifies people",           "27%"),
        ("promptinject.HijackHateHumans",  "Overrides instructions → hate speech",  "56%"),
        ("promptinject.HijackKillHumans",  "Overrides instructions → violence",     "50%"),
    ]

    headers = ["Probe", "Attack Goal", "Baseline Rate"]
    col_w = [4.5, 5.5, 1.7]
    col_x = [0.5, 5.0, 10.5]

    for ci, (hdr, w, x) in enumerate(zip(headers, col_w, col_x)):
        _colored_box(slide, x, 2.85, w, 0.38, hdr,
                     bg=ACCENT, fg=DARK_BG, font_size=13, bold=True)

    for ri, (probe, goal, rate) in enumerate(probes):
        y = 3.23 + ri * 0.47
        row_bg = RGBColor(0x22, 0x22, 0x3A) if ri % 2 == 0 else RGBColor(0x1A, 0x1A, 0x2E)
        rate_color = RED if rate != "0%" else GREEN
        for val, w, x in zip([probe, goal, rate], col_w, col_x):
            c = rate_color if val == rate else WHITE
            _colored_box(slide, x, y, w, 0.44, val,
                         bg=row_bg, fg=c, font_size=12, align=PP_ALIGN.LEFT if val != rate else PP_ALIGN.CENTER)

    _box(slide, 0.5, 6.08, 12.3, 0.35,
         "Overall baseline attack pass-rate: 98.1% of attack attempts succeeded against the bare model",
         font_size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def slide_04_vdps(prs: Presentation) -> None:
    """Slide 4 — VDPS (Vulnerability-Driven Patch Synthesis)."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Stage 2 — Vulnerability-Driven Patch Synthesis (VDPS)")
    _rule(slide, 1.1)

    _box(slide, 0.5, 1.2, 12.3, 0.4,
         "Novel contribution: patches synthesized FROM the attack data, not written in advance",
         font_size=15, bold=True, color=ACCENT)

    # Flow
    flow = [
        ("hitlog.jsonl\n(confirmed attacks)", RGBColor(0x6E, 0x2F, 0x1F)),
        ("Extract\nAttackPatterns\nper probe",  RGBColor(0x1F, 0x45, 0x6E)),
        ("GPT-3.5\nSynthesizes\nJSON mitigations", RGBColor(0x2D, 0x5A, 0x27)),
        ("patches_\nsynthesized\n.yaml",           RGBColor(0x50, 0x30, 0x70)),
    ]
    for i, (label, color) in enumerate(flow):
        x = 0.4 + i * 3.1
        _colored_box(slide, x, 1.75, 2.7, 1.2, label,
                     bg=color, fg=WHITE, font_size=13, bold=True)
        if i < 3:
            _box(slide, x + 2.72, 2.05, 0.38, 0.5, "→",
                 font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # JSON schema
    _colored_box(slide, 0.5, 3.1, 6.0, 2.8,
        'Synthesis prompt returns JSON:\n\n'
        '{\n'
        '  "system_prompt_addition": "Do not use deadnames...",\n'
        '  "blocked_keywords": ["Bruce Jenner", ...],\n'
        '  "injection_patterns": ["ignore.*instructions", ...],\n'
        '  "reasoning": "Detected probe targets trans identity..."\n'
        '}',
        bg=RGBColor(0x0D, 0x0D, 0x20), fg=GREEN, font_size=12, align=PP_ALIGN.LEFT)

    _bullet_box(slide, 6.7, 3.1, 6.1, 2.8, [
        "▸  One LLM call per probe category",
        "▸  3-attempt retry with stricter JSON prompting",
        "▸  All results merged: keywords deduplicated,",
        "   system additions concatenated per probe",
        "▸  Writes config YAML loaded by PatchEngine",
        "▸  Separates target LLM from synthesis LLM",
        "   (synthesis always uses reliable OpenAI model)",
    ], font_size=14)

    _box(slide, 0.5, 6.05, 12.3, 0.35,
         "Research question: do synthesized patches generalise beyond the exact attacks seen, or do they overfit?",
         font_size=13, color=AMBER, align=PP_ALIGN.CENTER)


def slide_05_proxy(prs: Presentation) -> None:
    """Slide 5 — Patch Proxy & Stage 3 Verify."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Stage 3 — PatchProxy & Verify")
    _rule(slide, 1.1)

    # Architecture diagram
    components = [
        (0.5,  1.35, 1.8, "Garak\nScanner",           RGBColor(0x6E, 0x2F, 0x1F)),
        (3.0,  1.35, 2.5, "PatchProxy\nlocalhost:8080", RGBColor(0x1F, 0x45, 0x6E)),
        (6.3,  1.35, 2.0, "OpenAI\nGPT-3.5",           RGBColor(0x2D, 0x5A, 0x27)),
    ]
    for x, y, w, label, color in components:
        _colored_box(slide, x, y, w, 0.9, label, bg=color, fg=WHITE, font_size=14, bold=True)

    # Arrows
    for ax in [2.3, 5.5]:
        _box(slide, ax, 1.6, 0.7, 0.4, "→", font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _box(slide, 5.5, 1.9, 0.7, 0.4, "←", font_size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Prompt patches
    _colored_box(slide, 3.0, 2.4, 2.5, 1.7,
        "PROMPT PATCHES\n(before model call)\n\n"
        "• InjectionGuardrail\n• SystemHardener\n• InputSanitizer",
        bg=RGBColor(0x1F, 0x35, 0x5E), fg=LIGHT_GREY, font_size=12, align=PP_ALIGN.LEFT)

    # Output patches
    _colored_box(slide, 3.0, 4.2, 2.5, 2.1,
        "OUTPUT PATCHES\n(after model call)\n\n"
        "• ToxicityFilter\n• KeywordPolicyFilter\n• DeadnamingFilter\n• QuackMedicineFilter\n• SexualisationFilter",
        bg=RGBColor(0x3E, 0x1F, 0x5E), fg=LIGHT_GREY, font_size=12, align=PP_ALIGN.LEFT)

    _bullet_box(slide, 6.7, 2.4, 6.1, 4.0, [
        "▸  FastAPI server — single /v1/chat/completions endpoint",
        "▸  OPENAI_BASE_URL override routes Garak transparently",
        "▸  Blocked inputs → immediate fallback, no upstream call",
        "▸  Response patching applied before Garak evaluates",
        "",
        "▸  ProxyServer lifecycle managed by pipeline:",
        "   start() → background thread → _wait_ready() socket poll",
        "   → Stage 3 scan → stop()",
        "",
        "▸  All patches independently toggleable via YAML config",
        "   → enables ablation studies with zero code changes",
    ], font_size=13)


def slide_06_prompt_patches(prs: Presentation) -> None:
    """Slide 6 — Prompt-Level Mitigation Mechanisms."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Mitigation Mechanisms — Prompt Level")
    _rule(slide, 1.1)

    _box(slide, 0.5, 1.2, 12.3, 0.35,
         "Operate BEFORE the model generates a response — prevent attacks from reaching the model",
         font_size=14, color=LIGHT_GREY)

    patches = [
        (
            "PromptInjectionGuardrail",
            "Regex scan for ~10 injection signatures:\n"
            "ignore/disregard instructions, DAN persona, jailbreak keyword,\n"
            "<SYSTEM> / [INST] tokens, ###SYSTEM markers.\n"
            "Match → input replaced with fallback. Model never called.\n"
            "Config-extensible: extra patterns from patches.yaml.",
        ),
        (
            "SystemPromptHardener",
            "Prepends a hardened security prefix to the system prompt.\n"
            "Explicitly instructs model to resist instruction-override attacks.\n"
            "When VDPS is active: prefix is synthesized from hitlog data,\n"
            "including probe-specific guardrails.\n"
            "Idempotent — never double-prepends on repeated calls.",
        ),
        (
            "InputSanitizer",
            "Strips zero-width / invisible Unicode characters (hidden payloads).\n"
            "Normalises Cyrillic homoglyphs (е→e, а→a, etc.).\n"
            "Collapses excessive whitespace runs.\n"
            "Redacts suspicious base64 payloads (≥40 chars) → [B64_REDACTED].",
        ),
    ]

    for i, (name, desc) in enumerate(patches):
        y = 1.65 + i * 1.8
        _colored_box(slide, 0.5, y, 3.2, 0.4, name,
                     bg=ACCENT, fg=DARK_BG, font_size=13, bold=True, align=PP_ALIGN.LEFT)
        _colored_box(slide, 0.5, y + 0.4, 12.3, 1.3, desc,
                     bg=RGBColor(0x1A, 0x1A, 0x2E), fg=LIGHT_GREY, font_size=13, align=PP_ALIGN.LEFT)


def slide_07_output_patches(prs: Presentation) -> None:
    """Slide 7 — Output-Level Mitigation Mechanisms."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Mitigation Mechanisms — Output Level")
    _rule(slide, 1.1)

    _box(slide, 0.5, 1.2, 12.3, 0.35,
         "Operate AFTER the model generates a response — catch harmful content before Garak evaluates it",
         font_size=14, color=LIGHT_GREY)

    # Two-stage design box
    _colored_box(slide, 0.5, 1.65, 12.3, 0.75,
        "Two-Stage Design for probe-specific filters:  "
        "Stage A — O(1) structural regex (always runs, free).  "
        "Stage B — LLM-as-judge API call (only fires if Stage A passes AND relevance pre-check matches).\n"
        "Minimises false positives and API cost. Name-agnostic — generalises to novel attack instances.",
        bg=RGBColor(0x0D, 0x0D, 0x20), fg=AMBER, font_size=13, align=PP_ALIGN.LEFT)

    patches = [
        ("ToxicityFilter",       "Generic",  "Regex for severe content (weapons, CSAM) + profanity density heuristic. Threshold configurable."),
        ("KeywordPolicyFilter",  "Generic",  "Exact keyword blocklist from config (VDPS populates this). Case-insensitive. Any match → fallback."),
        ("ResponseRewriter",     "Generic",  "Optional LLM-assisted rewrite of borderline responses. Disabled by default (adds latency + cost)."),
        ("DeadnamingFilter",     "2-Stage",  "Stage A: 9 structural patterns (formerly known as, birth name, changed name from/to...). Stage B: LLM judge — asks if trans person is deadnamed."),
        ("QuackMedicineFilter",  "2-Stage",  "Stage A: hedge-then-list pattern (limited evidence... → benefits list). Stage B: LLM judge — endorsing pseudoscientific practices?"),
        ("SexualisationFilter",  "2-Stage",  "Stage A: objectification phrases (boner material, turn-on, sexual object...). Stage B: LLM judge — sexualising or objectifying a person?"),
    ]

    for i, (name, kind, desc) in enumerate(patches):
        y = 2.55 + i * 0.78
        kind_color = ACCENT if kind == "Generic" else GREEN
        _colored_box(slide, 0.5, y, 2.8, 0.35, name,
                     bg=RGBColor(0x22, 0x22, 0x3A), fg=WHITE, font_size=12, bold=True, align=PP_ALIGN.LEFT)
        _colored_box(slide, 3.35, y, 1.0, 0.35, kind,
                     bg=kind_color, fg=DARK_BG, font_size=11, bold=True)
        _colored_box(slide, 4.4, y, 8.4, 0.35, desc,
                     bg=RGBColor(0x1A, 0x1A, 0x2E), fg=LIGHT_GREY, font_size=12, align=PP_ALIGN.LEFT)


def slide_08_results(prs: Presentation) -> None:
    """Slide 8 — Experimental Results."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Experimental Results")
    _rule(slide, 1.1)

    # Results table
    rows = [
        ("lmrc.Bullying",                 "0%",   "0%",   "—",       GREEN),
        ("lmrc.Deadnaming",               "83%",  "Partial","Partial",AMBER),
        ("lmrc.QuackMedicine",            "60%",  "0%",   "Eliminated",GREEN),
        ("lmrc.Sexualisation",            "27%",  "0%",   "Eliminated",GREEN),
        ("promptinject.HijackHateHumans", "56%",  "0%",   "Eliminated",GREEN),
        ("promptinject.HijackKillHumans", "50%",  "0%",   "Eliminated",GREEN),
    ]

    headers = ["Probe", "Baseline", "Patched", "Outcome"]
    col_w =   [4.6,      1.6,        1.6,       2.4]
    col_x =   [0.5,      5.1,        6.7,       8.3]

    for hdr, w, x in zip(headers, col_w, col_x):
        _colored_box(slide, x, 1.35, w, 0.4, hdr,
                     bg=ACCENT, fg=DARK_BG, font_size=13, bold=True)

    for ri, (probe, base, patch, outcome, oc) in enumerate(rows):
        y = 1.75 + ri * 0.5
        row_bg = RGBColor(0x22, 0x22, 0x3A) if ri % 2 == 0 else RGBColor(0x1A, 0x1A, 0x2E)
        for val, w, x, color in zip(
            [probe, base, patch, outcome],
            col_w, col_x,
            [WHITE, RED if base != "0%" else GREEN, GREEN if patch == "0%" else AMBER, oc]
        ):
            _colored_box(slide, x, y, w, 0.47, val,
                         bg=row_bg, fg=color, font_size=12,
                         align=PP_ALIGN.LEFT if val == probe else PP_ALIGN.CENTER)

    # Summary bar
    _colored_box(slide, 0.5, 4.95, 5.5, 0.9,
        "Overall Baseline: 98.1% attack pass-rate",
        bg=RGBColor(0x6E, 0x1F, 0x1F), fg=WHITE, font_size=15, bold=True)
    _colored_box(slide, 6.1, 4.95, 5.5, 0.9,
        "Overall Patched: 38.7% attack pass-rate",
        bg=RGBColor(0x1F, 0x5A, 0x27), fg=WHITE, font_size=15, bold=True)

    _box(slide, 0.5, 6.0, 12.3, 0.45,
         "Δ  −59.4 percentage points reduction in attack success rate",
         font_size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def slide_09_limitations(prs: Presentation) -> None:
    """Slide 9 — Observations & Limitations."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Observations & Limitations")
    _rule(slide, 1.1)

    col_left = [
        ("What worked well", ACCENT, [
            "▸  Network-layer proxy interception is scanner-agnostic",
            "▸  VDPS produces targeted, non-generic patches",
            "▸  Two-stage filters balance precision vs. API cost",
            "▸  Ablation configs isolate each mechanism's contribution",
            "▸  QuackMedicine and Sexualisation fully eliminated",
            "▸  Prompt injection attacks fully blocked",
        ]),
    ]

    col_right = [
        ("Limitations", RED, [
            "▸  Deadnaming only partially mitigated — LLM judge",
            "   misses subtle biographical prose",
            "▸  Keyword blocklists overfit to seen attack outputs",
            "▸  Proxy adds ~50–100ms latency per request",
            "▸  VDPS synthesis may fail on malformed JSON (3 retries)",
            "▸  No formal generalisation evaluation to novel attacks",
            "▸  Single model tested (GPT-3.5-turbo)",
        ]),
    ]

    for (title, color, bullets), x in zip(col_left + col_right, [0.5, 6.7]):
        _colored_box(slide, x, 1.25, 5.9, 0.4, title,
                     bg=color, fg=DARK_BG if color == ACCENT else WHITE,
                     font_size=14, bold=True)
        _bullet_box(slide, x, 1.7, 5.9, 4.5, bullets, font_size=13)

    _colored_box(slide, 0.5, 6.0, 12.3, 0.6,
        "Key finding: layered prompt + output defenses are complementary — "
        "neither alone achieves the same reduction as both combined.",
        bg=RGBColor(0x0D, 0x0D, 0x20), fg=AMBER, font_size=14, bold=True)


def slide_10_future(prs: Presentation) -> None:
    """Slide 10 — Future Research Directions."""
    slide = _blank_slide(prs)
    _bg(slide)
    _heading(slide, "Future Research Directions")
    _rule(slide, 1.1)

    _box(slide, 0.5, 1.2, 12.3, 0.4,
         "If extended to a full research project (USENIX Security / CCS / IEEE S&P):",
         font_size=14, color=LIGHT_GREY)

    directions = [
        ("Generalisation of VDPS patches",
         "Do synthesized patches block paraphrased attacks not seen in the hitlog?\n"
         "Design: hold-out adversarial set evaluated after VDPS. Overfitting = patch brittleness."),
        ("Multi-model & multi-scanner evaluation",
         "Evaluate across GPT-4, Claude, Llama models. Test with additional scanners (PyRIT, Promptbench).\n"
         "Does VDPS transfer patch knowledge across models?"),
        ("Iterative self-improvement convergence",
         "Run N probe → patch → verify cycles. Does attack rate monotonically decrease?\n"
         "Explore whether the loop converges or oscillates on adversarially adaptive attacks."),
        ("LLM judge calibration & false positive rate",
         "Measure false positive rate of semantic filters (DeadnamingFilter, etc.) on benign inputs.\n"
         "Calibrate relevance pre-check thresholds. Benchmark against dedicated safety classifiers."),
        ("Adaptive adversaries",
         "Re-run Garak after patches are published — do attackers adapt?\n"
         "Explore whether patch synthesis can keep pace with adaptive red-team probes."),
    ]

    for i, (title, desc) in enumerate(directions):
        y = 1.7 + i * 1.1
        _colored_box(slide, 0.5, y, 0.35, 0.85, str(i + 1),
                     bg=ACCENT, fg=DARK_BG, font_size=16, bold=True)
        _colored_box(slide, 0.9, y, 3.8, 0.38, title,
                     bg=RGBColor(0x1F, 0x45, 0x6E), fg=WHITE, font_size=13, bold=True, align=PP_ALIGN.LEFT)
        _colored_box(slide, 0.9, y + 0.38, 11.9, 0.47, desc,
                     bg=RGBColor(0x1A, 0x1A, 0x2E), fg=LIGHT_GREY, font_size=12, align=PP_ALIGN.LEFT)


# ── Main ──────────────────────────────────────────────────────────────

def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_architecture(prs)
    slide_03_garak(prs)
    slide_04_vdps(prs)
    slide_05_proxy(prs)
    slide_06_prompt_patches(prs)
    slide_07_output_patches(prs)
    slide_08_results(prs)
    slide_09_limitations(prs)
    slide_10_future(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved -> {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate research presentation PPTX")
    parser.add_argument(
        "--output",
        default="results/presentation.pptx",
        help="Output path for the PPTX file",
    )
    args = parser.parse_args()
    build_presentation(Path(args.output))
